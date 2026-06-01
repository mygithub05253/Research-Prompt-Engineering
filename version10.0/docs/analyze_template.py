# GIC 양식 ver2.pptx 분석 스크립트 (일회성)
# 용도: 슬라이드별 텍스트 박스 위치·크기·폰트·컬러 추출
# 결과: design/양식_분석_결과.md 와 design/양식_분석_결과.json

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor

BASE = Path(r"C:\Users\kik32\내 드라이브\Obsidian Vault\Obsidian_School\가천대학교\GIC\4기\GIC 리서치 프롬프트 모델 생성 프로젝트\version10.0")
SRC = BASE / "design" / "GIC_양식_ver2.pptx"
OUT_MD = BASE / "design" / "양식_분석_결과.md"
OUT_JSON = BASE / "design" / "양식_분석_결과.json"


def emu_to_inch(emu):
    return round(Emu(emu).inches, 3) if emu is not None else None


def emu_to_cm(emu):
    return round(Emu(emu).cm, 2) if emu is not None else None


def color_to_hex(color_format):
    try:
        if color_format and color_format.type is not None:
            rgb = color_format.rgb
            if rgb:
                return f"#{rgb}"
    except Exception:
        pass
    return None


def extract_run_info(run):
    info = {
        "text": run.text,
        "font_name": None,
        "font_size_pt": None,
        "bold": None,
        "color": None,
    }
    try:
        info["font_name"] = run.font.name
    except Exception:
        pass
    try:
        if run.font.size:
            info["font_size_pt"] = run.font.size.pt
    except Exception:
        pass
    try:
        info["bold"] = run.font.bold
    except Exception:
        pass
    try:
        info["color"] = color_to_hex(run.font.color)
    except Exception:
        pass
    return info


def extract_shape_info(shape, slide_idx):
    base = {
        "slide_idx": slide_idx,
        "shape_id": shape.shape_id,
        "shape_name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type else None,
        "left_in": emu_to_inch(shape.left),
        "top_in": emu_to_inch(shape.top),
        "width_in": emu_to_inch(shape.width),
        "height_in": emu_to_inch(shape.height),
        "left_cm": emu_to_cm(shape.left),
        "top_cm": emu_to_cm(shape.top),
        "width_cm": emu_to_cm(shape.width),
        "height_cm": emu_to_cm(shape.height),
    }

    fill_hex = None
    try:
        if shape.fill.type is not None:
            fill_hex = color_to_hex(shape.fill.fore_color)
    except Exception:
        pass
    base["fill_color"] = fill_hex

    paragraphs = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            runs = [extract_run_info(r) for r in para.runs]
            text = "".join(r["text"] for r in runs)
            if text.strip() or runs:
                paragraphs.append({"text": text, "runs": runs, "alignment": str(para.alignment) if para.alignment else None})
    base["paragraphs"] = paragraphs
    base["text_total"] = "\n".join(p["text"] for p in paragraphs)
    base["char_count"] = len(base["text_total"])
    return base


def main():
    prs = Presentation(SRC)
    slide_w_in = emu_to_inch(prs.slide_width)
    slide_h_in = emu_to_inch(prs.slide_height)
    slide_w_cm = emu_to_cm(prs.slide_width)
    slide_h_cm = emu_to_cm(prs.slide_height)

    result = {
        "source": str(SRC),
        "slide_count": len(prs.slides),
        "slide_size": {"w_in": slide_w_in, "h_in": slide_h_in, "w_cm": slide_w_cm, "h_cm": slide_h_cm},
        "slides": [],
        "color_palette": set(),
        "font_palette": set(),
    }

    for idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else None
        shapes_info = []
        for shape in slide.shapes:
            sinfo = extract_shape_info(shape, idx)
            shapes_info.append(sinfo)
            if sinfo["fill_color"]:
                result["color_palette"].add(sinfo["fill_color"])
            for para in sinfo["paragraphs"]:
                for run in para["runs"]:
                    if run["color"]:
                        result["color_palette"].add(run["color"])
                    if run["font_name"]:
                        result["font_palette"].add(run["font_name"])
        result["slides"].append({
            "idx": idx,
            "layout_name": layout_name,
            "shape_count": len(shapes_info),
            "shapes": shapes_info,
        })

    result["color_palette"] = sorted(result["color_palette"])
    result["font_palette"] = sorted(result["font_palette"])

    OUT_JSON.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")

    md = []
    md.append(f"# GIC 양식 ver2.pptx 분석 결과")
    md.append("")
    md.append(f"- **출처**: `{SRC.name}`")
    md.append(f"- **슬라이드 수**: {result['slide_count']}")
    md.append(f"- **슬라이드 크기**: {slide_w_in}″ × {slide_h_in}″ ({slide_w_cm}cm × {slide_h_cm}cm)")
    md.append("")
    md.append("## 컬러 팔레트")
    for c in result["color_palette"]:
        md.append(f"- `{c}`")
    md.append("")
    md.append("## 폰트 팔레트")
    for f in result["font_palette"]:
        md.append(f"- `{f}`")
    md.append("")
    md.append("## 슬라이드별 도형 상세")
    for s in result["slides"]:
        md.append(f"\n### 슬라이드 {s['idx']+1} — 레이아웃: `{s['layout_name']}` (도형 {s['shape_count']}개)")
        md.append("\n| # | 이름 | 위치(cm) | 크기(cm) | 텍스트(앞 60자) | 채움 |")
        md.append("|---|---|---|---|---|---|")
        for i, sh in enumerate(s["shapes"]):
            txt = (sh["text_total"] or "").replace("\n", " / ")[:60]
            pos = f"L{sh['left_cm']} T{sh['top_cm']}" if sh["left_cm"] is not None else "-"
            size = f"W{sh['width_cm']} H{sh['height_cm']}" if sh["width_cm"] is not None else "-"
            md.append(f"| {i+1} | {sh['shape_name']} | {pos} | {size} | {txt} | {sh['fill_color'] or '-'} |")

    OUT_MD.write_text("\n".join(md), encoding="utf-8")
    print(f"OK -> {OUT_MD}")
    print(f"OK -> {OUT_JSON}")
    print(f"slides: {result['slide_count']} / colors: {len(result['color_palette'])} / fonts: {len(result['font_palette'])}")


if __name__ == "__main__":
    main()
