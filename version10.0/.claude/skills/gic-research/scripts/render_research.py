# GIC 기업리서치 PPT 자동 생성 (v10.0 트랙 1)
# 입력: 부록 C yaml + templates/기업리서치_9p.pptx (학회 양식 ver2 복제)
# 출력: data/output/[기업명]_GIC리서치_YYYYMMDD/[기업명]_GIC리서치_YYYYMMDD.pptx
#
# v0.1 — 텍스트 박스(종목명·투자의견·Point 1·2·3·우상단 기업명·작성일·학회기수·Check Point) 자동 매핑.
# 표·이미지 영역은 양식.pptx 그대로 유지 (부원이 PowerPoint에서 수기 채움).
# v0.2(예정) — 표 행 자동 채움(Revision/시총/주요주주/재무지표) + 이미지 영역 PNG 삽입.

from __future__ import annotations
import argparse
import io
import sys
from datetime import datetime
from pathlib import Path

# Windows cp949 환경에서 유니코드 출력 보장
if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

import yaml
from pptx import Presentation
from pptx.util import Pt


# ─────────────────────────────────────────────────────────────
# 유틸: 텍스트 박스 안의 문자열 안전 교체 (run 단위)
# ─────────────────────────────────────────────────────────────
def replace_in_textframe(tf, old: str, new: str) -> bool:
    """run.text 단위로 안전하게 교체. 매칭 실패 시 paragraph 단위 재구성(포맷 일부 손실 가능)."""
    if not tf or not old:
        return False
    # 1차: run 단위 정확 매칭
    for para in tf.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True
    # 2차: paragraph 단위 (run이 분리됐을 때)
    for para in tf.paragraphs:
        if old in para.text:
            full = para.text.replace(old, new)
            # 첫 run에 새 텍스트, 나머지 run 비움
            if para.runs:
                for i, run in enumerate(para.runs):
                    run.text = full if i == 0 else ""
                return True
    return False


def set_textframe(tf, text: str) -> None:
    """텍스트 박스를 단일 텍스트로 덮어쓰기 (포맷은 첫 run에서 상속)."""
    if not tf:
        return
    if tf.paragraphs:
        first_para = tf.paragraphs[0]
        if first_para.runs:
            first_run = first_para.runs[0]
            first_run.text = text
            # 나머지 run/paragraph는 비움
            for run in list(first_para.runs)[1:]:
                run.text = ""
            for para in list(tf.paragraphs)[1:]:
                for run in para.runs:
                    run.text = ""
            return
    tf.text = text


# ─────────────────────────────────────────────────────────────
# 슬라이드 1 — 표지 처리
# ─────────────────────────────────────────────────────────────
def render_slide_1_cover(slide, data: dict) -> list[str]:
    log = []
    s = data.get("slide_1_cover", {})

    종목명 = str(s.get("종목명", "[기업명]"))
    종목코드 = str(s.get("종목코드", "[######]"))
    부제 = str(s.get("부제", "GIC여 긱스럽게 도전하라"))
    작성일 = str(s.get("작성일", datetime.now().strftime("%Y.%m.%d")))
    학회기수 = str(s.get("학회기수", "GIC 4기"))
    투자의견 = str(s.get("투자의견", "BUY"))
    투자의견_한글 = str(s.get("투자의견_한글", "매수"))
    check_point_슬로건 = str(s.get("check_point", {}).get("슬로건", ""))

    point_1 = s.get("point_1", {}) or {}
    point_2 = s.get("point_2", {}) or {}
    point_3 = s.get("point_3", {}) or {}

    # 좌상단 작성자 자리(TextBox 7) — 학회 양식 ver2 placeholder "GIC 1기 최진호"
    # 우상단 학회기수(TextBox 21) — placeholder "GIC 3기 가나다"
    # 학회기수 yaml 1개로 두 자리 모두 갱신
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        text = tf.text

        # 종목명 + 종목코드 : "지아이씨 (240905)"
        if "지아이씨" in text or "240905" in text:
            replace_in_textframe(tf, "지아이씨", 종목명)
            replace_in_textframe(tf, "240905", 종목코드)
            log.append(f"slide_1.종목명/종목코드: {종목명} ({종목코드})")
            continue

        # 부제 : "GIC여 긱스럽게 도전하라" — 그대로 유지(yaml에서 변경 가능)
        if "긱스럽게" in text and 부제 != "GIC여 긱스럽게 도전하라":
            replace_in_textframe(tf, "GIC여 긱스럽게 도전하라", 부제)
            log.append(f"slide_1.부제: {부제}")
            continue

        # 작성일: "2026.01.01"
        if "2026.01.01" in text:
            replace_in_textframe(tf, "2026.01.01", 작성일)
            log.append(f"slide_1.작성일: {작성일}")
            continue

        # 학회기수 우상단: "GIC 3기 가나다"
        if "GIC 3기 가나다" in text:
            replace_in_textframe(tf, "GIC 3기 가나다", 학회기수)
            log.append(f"slide_1.학회기수: {학회기수}")
            continue

        # 좌상단 작성자: "GIC 1기 최진호"
        if "GIC 1기 최진호" in text:
            replace_in_textframe(tf, "GIC 1기 최진호", 학회기수)
            log.append(f"slide_1.작성자: {학회기수}")
            continue

        # 투자의견 박스 " BUY(매수)" — 양식 placeholder
        if "BUY(매수)" in text or "BUY" in text:
            # 영문 + 한글 동시 교체
            replace_in_textframe(tf, "BUY", 투자의견)
            replace_in_textframe(tf, "매수", 투자의견_한글)
            log.append(f"slide_1.투자의견: {투자의견}({투자의견_한글})")
            continue

        # Check Point 슬로건
        if "내가 그린 기린 그림은 잘 그린 기린 그림" in text:
            if check_point_슬로건:
                replace_in_textframe(
                    tf, "내가 그린 기린 그림은 잘 그린 기린 그림", check_point_슬로건
                )
                log.append(f"slide_1.check_point.슬로건: {check_point_슬로건}")
            continue

        # Point 1
        if text.strip().startswith("Point 1."):
            제목 = str(point_1.get("제목", "") or "")
            내용 = str(point_1.get("내용", "") or "")
            if 제목 or 내용:
                # 첫 줄: "Point 1. [제목]" / 둘째 줄부터: 내용
                set_textframe(tf, f"Point 1. {제목}\n{내용}".strip())
                log.append(f"slide_1.point_1: {제목[:20]}…")
            continue
        if text.strip().startswith("Point 2."):
            제목 = str(point_2.get("제목", "") or "")
            내용 = str(point_2.get("내용", "") or "")
            if 제목 or 내용:
                set_textframe(tf, f"Point 2. {제목}\n{내용}".strip())
                log.append(f"slide_1.point_2: {제목[:20]}…")
            continue
        if text.strip().startswith("Point 3."):
            제목 = str(point_3.get("제목", "") or "")
            내용 = str(point_3.get("내용", "") or "")
            if 제목 or 내용:
                set_textframe(tf, f"Point 3. {제목}\n{내용}".strip())
                log.append(f"slide_1.point_3: {제목[:20]}…")
            continue

    return log


# ─────────────────────────────────────────────────────────────
# 슬라이드 2~9 — 우상단 기업명 라벨 (공통)
# ─────────────────────────────────────────────────────────────
def render_company_label(slide, 기업명_라벨: str) -> bool:
    """우상단 '기업명(번호)' placeholder를 실제 라벨로 교체"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if "기업명(번호)" in tf.text:
            replace_in_textframe(tf, "기업명(번호)", 기업명_라벨)
            return True
    return False


# ─────────────────────────────────────────────────────────────
# word_wrap=False 후처리 (design-guide_v10 §11.9)
# ─────────────────────────────────────────────────────────────
def apply_word_wrap_postprocess(prs) -> int:
    """짧은 라벨 (≤30자, 1 paragraph)에 word_wrap=False"""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text.strip()
            if len(text) <= 30 and len(tf.paragraphs) == 1:
                tf.word_wrap = False
                count += 1
    return count


# ─────────────────────────────────────────────────────────────
# 메인 렌더
# ─────────────────────────────────────────────────────────────
def render(yaml_path: Path, template_path: Path, output_path: Path) -> dict:
    if not yaml_path.exists():
        raise FileNotFoundError(f"yaml 파일 없음: {yaml_path}")
    if not template_path.exists():
        raise FileNotFoundError(f"양식 템플릿 없음: {template_path}")

    data = yaml.safe_load(yaml_path.read_text(encoding="utf-8"))
    if not data:
        raise ValueError("yaml 파싱 결과 비어있음")

    prs = Presentation(str(template_path))

    log_lines = []
    log_lines.append(f"=== GIC 리서치 PPT 렌더 로그 ({datetime.now().isoformat()}) ===")
    log_lines.append(f"입력 yaml: {yaml_path}")
    log_lines.append(f"양식 템플릿: {template_path}")
    log_lines.append(f"슬라이드 수: {len(prs.slides)}")
    log_lines.append("")

    # 슬라이드 1
    log_lines.append("[슬라이드 1 — 표지]")
    log_lines.extend("  " + l for l in render_slide_1_cover(prs.slides[0], data))
    log_lines.append("")

    # 슬라이드 2~9 — 우상단 기업명
    s1 = data.get("slide_1_cover", {})
    종목명 = str(s1.get("종목명", "[기업명]"))
    종목코드 = str(s1.get("종목코드", "[######]"))
    기업명_라벨 = f"{종목명}({종목코드})"

    log_lines.append("[슬라이드 2~9 — 공통 우상단 라벨]")
    for i in range(1, len(prs.slides)):
        ok = render_company_label(prs.slides[i], 기업명_라벨)
        log_lines.append(f"  슬라이드 {i+1}: {'OK' if ok else 'placeholder 없음'}")
    log_lines.append("")

    # 후처리
    wrap_count = apply_word_wrap_postprocess(prs)
    log_lines.append(f"[후처리] word_wrap=False 적용: {wrap_count}건")
    log_lines.append("")

    # 저장
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    log_lines.append(f"저장 완료: {output_path}")

    log_text = "\n".join(log_lines)
    log_path = output_path.parent / "render_log.txt"
    log_path.write_text(log_text, encoding="utf-8")

    return {
        "output_pptx": output_path,
        "log_path": log_path,
        "log_text": log_text,
        "slide_count": len(prs.slides),
    }


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--yaml", required=True, help="부록 C yaml 파일 경로")
    ap.add_argument("--template", default=None, help="양식 템플릿 경로 (기본: ../../../../templates/기업리서치_9p.pptx)")
    ap.add_argument("--output", required=True, help="출력 PPTX 경로")
    args = ap.parse_args()

    yaml_path = Path(args.yaml)
    if args.template:
        template_path = Path(args.template)
    else:
        # 기본: render_research.py에서 templates/기업리서치_9p.pptx까지 4단계 상위
        # .claude/skills/gic-research/scripts/render_research.py → version10.0/templates/...
        script = Path(__file__).resolve()
        v10 = script.parents[4]
        template_path = v10 / "templates" / "기업리서치_9p.pptx"
    output_path = Path(args.output)

    try:
        result = render(yaml_path, template_path, output_path)
        print(result["log_text"])
        sys.exit(0)
    except Exception as e:
        print(f"ERROR: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
